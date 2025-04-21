import pdfplumber
from docx import Document
import re
from fastapi import UploadFile
from pptx import Presentation

def extract_pdf_text(file: UploadFile):
    with pdfplumber.open(file.file) as pdf:
        text = ""
        for page in pdf.pages:
            text += page.extract_text()
    return text

def extract_docx_text(file: UploadFile):
    doc = Document(file.file)
    text = ""
    for para in doc.paragraphs:
        text += para.text
    return text

def extract_pptx_text(file: UploadFile):
    presentation = Presentation(file.file)
    text = ""
    
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def extract_txt_text(file: UploadFile):
    contents = file.file.read().decode("utf-8")
    return contents